"""
build_sales_pptx.py
───────────────────
AI Trade Management — 営業スライド（初回商談用）生成スクリプト
15スライド構成: 表紙→アジェンダ→Why Now→AMAT事例→As-Is→課題→To-Be→Gap→Value Prop→製品→競合優位→実績→ロードマップ→ROI→CTA

Usage:
    python scripts/build_sales_pptx.py
    python scripts/build_sales_pptx.py --output data/demo/my_deck.pptx
"""

from __future__ import annotations
import argparse
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── 出力先 ──────────────────────────────────────────────────
ROOT = Path(__file__).parent.parent
OUT_DIR = ROOT / "data" / "demo"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# ── スライドサイズ (16:9) ────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── カラーパレット ────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x08, 0x0d, 0x1a)   # #080d1a  ダーク背景
C_BG_NAVY   = RGBColor(0x0f, 0x16, 0x27)   # #0f1627  ミッドナビー
C_BG_LIGHT  = RGBColor(0xf8, 0xfa, 0xff)   # #f8faff  ライト背景
C_BLUE      = RGBColor(0x0a, 0x6c, 0xff)   # #0a6cff  プライマリブルー
C_BLUE_LT   = RGBColor(0x60, 0xa5, 0xfa)   # #60a5fa  ライトブルー
C_TEXT_W    = RGBColor(0xf0, 0xf4, 0xff)   # #f0f4ff  テキスト白
C_TEXT_D    = RGBColor(0x0d, 0x12, 0x20)   # #0d1220  テキストダーク
C_GRAY      = RGBColor(0x88, 0x92, 0xa4)   # #8892a4  グレー
C_CARD_D    = RGBColor(0x0d, 0x12, 0x20)   # #0d1220  ダークカード
C_BORDER    = RGBColor(0x1e, 0x2a, 0x3a)   # #1e2a3a  ボーダー
C_RED       = RGBColor(0xf7, 0x44, 0x44)   # #f74444  レッド
C_GREEN     = RGBColor(0x4a, 0xde, 0x80)   # #4ade80  グリーン
C_GOLD      = RGBColor(0xc8, 0xa9, 0x6e)   # #c8a96e  ゴールド
C_INDIGO    = RGBColor(0x81, 0x8c, 0xf8)   # #818cf8  インディゴ
C_WARN_BG   = RGBColor(0x11, 0x0a, 0x0a)   # #110a0a  警告カード背景
C_GOOD_BG   = RGBColor(0x06, 0x0e, 0x1a)   # #060e1a  良好カード背景
C_WHITE     = RGBColor(0xff, 0xff, 0xff)
C_LIGHT_BORDER = RGBColor(0xd8, 0xe4, 0xf0)

# ── ヘルパー関数 ─────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_rect(
    slide, x, y, w, h,
    fill: RGBColor | None = None,
    line_color: RGBColor | None = None,
    line_width_pt: float = 0.75,
    alpha: int | None = None,
) -> object:
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = Pt(0) if line_color is None else Pt(line_width_pt)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide, text: str,
    x, y, w, h,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = C_TEXT_W,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
    font_name: str = "Noto Sans JP",
) -> object:
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_label_title_body(
    slide,
    label: str, title: str, body: str,
    x, y, w,
    is_dark: bool = True,
    accent_color: RGBColor = None,
) -> None:
    """ラベル + タイトル + 本文の縦積みテキストブロック"""
    ac = accent_color or (C_BLUE if is_dark else C_BLUE)
    tc = C_TEXT_W if is_dark else C_TEXT_D
    bc = C_GRAY if is_dark else RGBColor(0x4a, 0x55, 0x68)

    # accent bar
    add_rect(slide, x, y, Inches(0.07), Inches(0.6), fill=ac)

    lx = x + Inches(0.14)
    lw = w - Inches(0.14)
    add_textbox(slide, label, lx, y, lw, Inches(0.22),
                font_size=8, bold=True, color=ac, font_name="Arial")
    add_textbox(slide, title, lx, y + Inches(0.2), lw, Inches(0.5),
                font_size=20, bold=True, color=tc)
    add_textbox(slide, body, lx, y + Inches(0.65), lw, Inches(0.6),
                font_size=10, color=bc)


def add_footer(slide, is_dark: bool = True, slide_num: str = "01 / 15") -> None:
    brand = "AI TRADE COMPLIANCE MANAGEMENT"
    c = C_BORDER if is_dark else C_LIGHT_BORDER
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.33), Pt(1), fill=c)
    tc = RGBColor(0x2d, 0x37, 0x48) if is_dark else RGBColor(0xb8, 0xc4, 0xd4)
    add_textbox(slide, brand, Inches(0.5), Inches(7.1), Inches(8), Inches(0.25),
                font_size=7, bold=True, color=tc, font_name="Arial")
    add_textbox(slide, slide_num, Inches(11.8), Inches(7.1), Inches(1), Inches(0.25),
                font_size=9, bold=True, color=tc, align=PP_ALIGN.RIGHT, font_name="Arial")


def card_dark(slide, x, y, w, h, icon: str, tag: str, title: str, body: str,
              border_color: RGBColor = None, bg: RGBColor = None) -> None:
    bc = border_color or C_BORDER
    bg_c = bg or C_CARD_D
    add_rect(slide, x, y, w, h, fill=bg_c, line_color=bc, line_width_pt=0.75)
    cx = x + Inches(0.15)
    cw = w - Inches(0.3)
    add_textbox(slide, icon, cx, y + Inches(0.08), Inches(0.35), Inches(0.35), font_size=16, color=C_TEXT_W)
    add_textbox(slide, tag,  cx + Inches(0.38), y + Inches(0.1), cw - Inches(0.38), Inches(0.22),
                font_size=8, bold=True, color=C_BLUE, font_name="Arial")
    add_textbox(slide, title, cx, y + Inches(0.38), cw, Inches(0.3),
                font_size=11, bold=True, color=C_TEXT_W)
    add_textbox(slide, body, cx, y + Inches(0.66), cw, h - Inches(0.78),
                font_size=9, color=C_GRAY)


def card_light(slide, x, y, w, h, icon: str, tag: str, title: str, body: str,
               border_color: RGBColor = None, bg: RGBColor = None) -> None:
    bc = border_color or C_LIGHT_BORDER
    bg_c = bg or C_WHITE
    add_rect(slide, x, y, w, h, fill=bg_c, line_color=bc, line_width_pt=0.75)
    cx = x + Inches(0.15)
    cw = w - Inches(0.3)
    add_textbox(slide, icon, cx, y + Inches(0.08), Inches(0.35), Inches(0.35), font_size=16, color=C_TEXT_D)
    add_textbox(slide, tag,  cx + Inches(0.38), y + Inches(0.1), cw - Inches(0.38), Inches(0.22),
                font_size=8, bold=True, color=C_BLUE, font_name="Arial")
    add_textbox(slide, title, cx, y + Inches(0.38), cw, Inches(0.3),
                font_size=11, bold=True, color=C_TEXT_D)
    add_textbox(slide, body, cx, y + Inches(0.66), cw, h - Inches(0.78),
                font_size=9, color=RGBColor(0x4a, 0x55, 0x68))


def stat_box(slide, x, y, w, h, number: str, label: str,
             num_color: RGBColor = None, is_dark: bool = True) -> None:
    nc = num_color or C_BLUE
    bg = C_CARD_D if is_dark else C_WHITE
    bc = C_BORDER if is_dark else C_LIGHT_BORDER
    lc = C_GRAY if is_dark else RGBColor(0x4a, 0x55, 0x68)
    add_rect(slide, x, y, w, h, fill=bg, line_color=bc, line_width_pt=0.75)
    add_textbox(slide, number, x, y + Inches(0.1), w, Inches(0.45),
                font_size=22, bold=True, color=nc, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x + Inches(0.05), y + Inches(0.52), w - Inches(0.1), h - Inches(0.55),
                font_size=8, color=lc, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ════════════════════════════════════════════════════════════

def slide_01_cover(prs: Presentation) -> None:
    """表紙 — DARK"""
    s = blank_slide(prs)
    # Background
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_DARK)
    # Grid lines (subtle blue)
    for i in range(0, 14):
        c = RGBColor(0x0a, 0x6c, 0xff)
        r = add_rect(s, Inches(i * 1), 0, Pt(0.5), SLIDE_H, fill=c)
        r.fill.fore_color.rgb = c
        # set low transparency via alpha hack
    # Accent bar left
    add_rect(s, Inches(0.5), Inches(1.8), Inches(0.08), Inches(1.5), fill=C_BLUE)

    # Label
    add_textbox(s, "FIRST MEETING  |  CONFIDENTIAL",
                Inches(0.72), Inches(1.8), Inches(8), Inches(0.25),
                font_size=8, bold=True, color=C_BLUE, font_name="Arial")

    # Main title
    add_textbox(s, "輸出管理を、",
                Inches(0.72), Inches(2.15), Inches(10), Inches(0.7),
                font_size=36, bold=True, color=C_TEXT_W)
    add_textbox(s, "AIと法の力で再定義する。",
                Inches(0.72), Inches(2.75), Inches(11), Inches(0.7),
                font_size=36, bold=True, color=C_BLUE_LT)

    # Sub text
    add_textbox(s,
                "NeuroSymbolicアーキテクチャが実現する、説明できるコンプライアンス判断プラットフォーム",
                Inches(0.86), Inches(3.55), Inches(10), Inches(0.5),
                font_size=12, color=C_GRAY)

    # Product name
    add_textbox(s, "AI Trade Compliance Management",
                Inches(0.72), Inches(4.3), Inches(8), Inches(0.4),
                font_size=11, bold=True, color=C_BLUE, font_name="Arial")

    # Divider
    add_rect(s, Inches(0.5), Inches(6.5), Inches(12.33), Pt(1), fill=C_BORDER)

    # Bottom info
    add_textbox(s, "DATE: 2026.03    |    VERSION: v2.0    |    CONFIDENTIAL",
                Inches(0.5), Inches(6.6), Inches(8), Inches(0.3),
                font_size=8, color=C_GRAY, font_name="Arial")

    add_footer(s, is_dark=True, slide_num="01 / 15")


def slide_02_agenda(prs: Presentation) -> None:
    """アジェンダ — DARK"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_DARK)

    add_label_title_body(
        s, "TODAY'S AGENDA", "本日30分でお伝えすること",
        "規制環境の現実から始め、課題の根本に触れ、解決策の核心をお伝えします",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=True,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1),
             fill=C_BLUE)

    items = [
        ("01", "規制環境の激変 — なぜ今が転換点か",  "Entity List 3倍増・制裁総額2.1兆円の現実",        C_RED),
        ("02", "As-Is / Gap / To-Be — 現状体制と変革の必要性",  "属人化・コスト・キャッチオール——3つの根本課題",    C_GOLD),
        ("03", "AI Trade Management — 私たちが解決する方法",    "NeuroSymbolicアーキテクチャ・8モジュール・実証データ", C_BLUE),
        ("04", "導入ロードマップ・ROI・次のステップ",          "30日無償パイロット → デモ → フィット&ギャップ",     C_GREEN),
    ]
    ay = Inches(1.5)
    item_h = Inches(1.15)
    for num, main, sub, color in items:
        add_rect(s, Inches(0.5), ay, Inches(12.33), item_h - Inches(0.08),
                 fill=C_CARD_D, line_color=C_BORDER, line_width_pt=0.75)
        add_textbox(s, num, Inches(0.7), ay + Inches(0.25), Inches(0.5), Inches(0.5),
                    font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Arial")
        add_textbox(s, main, Inches(1.35), ay + Inches(0.12), Inches(10.0), Inches(0.35),
                    font_size=13, bold=True, color=C_TEXT_W)
        add_textbox(s, sub, Inches(1.35), ay + Inches(0.48), Inches(10.0), Inches(0.3),
                    font_size=10, color=C_GRAY)
        ay += item_h

    add_footer(s, is_dark=True, slide_num="02 / 15")


def slide_03_why_now(prs: Presentation) -> None:
    """Why Now: 規制環境の転換点 — NAVY"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_NAVY)

    add_label_title_body(
        s, "WHY NOW — MARKET CONTEXT",
        "規制環境の転換点——今が待ったなしの理由",
        "地政学リスクの高まりとともに、輸出規制の複雑性は急速に増している",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=True,
        accent_color=C_GOLD,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_GOLD)

    # Stat row
    stats = [
        ("3×",    "Entity List掲載数\n（2019→2024年）",   C_RED),
        ("2.1兆円","BIS違反制裁総額\n（2020〜2024累計）", C_RED),
        ("16件/年","日本企業が受けた\n外為法行政指導",     C_GOLD),
        ("78%",   "「現行体制では限界」\n輸出管理担当者", C_RED),
    ]
    sw = Inches(2.9)
    sx = Inches(0.5)
    for number, label, color in stats:
        stat_box(s, sx, Inches(1.45), sw, Inches(1.0), number, label, num_color=color, is_dark=True)
        sx += Inches(3.05)

    # 3 cards
    cards = [
        ("🌏", "地政学", "米中技術覇権競争の激化",
         "半導体・AI・量子分野を中心に米国の輸出管理強化が加速。2022年FDP規則改定で日本企業も規制対象に直接影響。"),
        ("📋", "規制複雑化", "外為法・EAR・ECCNの三重管理",
         "日本の外為法（輸出令・外為令）＋米国EAR（ECCN分類）＋HSコードの三重管理が常態化。一取引に複数規制が交差。"),
        ("⚡", "速度乖離", "ビジネス速度と規制対応の乖離",
         "グローバル商談はリアルタイム化。規制判断は依然として数週間単位。この速度差が商機損失と「判断しないリスク」を生む。"),
    ]
    cw = Inches(4.0)
    cx = Inches(0.5)
    cy = Inches(2.6)
    ch = Inches(3.8)
    for icon, tag, title, body in cards:
        card_dark(s, cx, cy, cw, ch, icon, tag, title, body,
                  border_color=RGBColor(0x1e, 0x3a, 0x2a) if tag == "地政学" else C_BORDER,
                  bg=C_WARN_BG if tag != "規制複雑化" else C_CARD_D)
        cx += Inches(4.17)

    add_footer(s, is_dark=True, slide_num="03 / 15")


def slide_04_amat(prs: Presentation) -> None:
    """AMAT制裁事例 — RED DARK"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_DARK)
    # Red top bar
    add_rect(s, 0, 0, SLIDE_W, Inches(0.06), fill=C_RED)

    add_label_title_body(
        s, "CASE STUDY — 制裁実例",
        "$153M超の制裁金——AMATの教訓",
        "これは「他人事」ではない。規制の複雑性は企業規模を問わない。",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=True,
        accent_color=C_RED,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_RED)

    # Left column
    add_rect(s, Inches(0.5), Inches(1.45), Inches(5.5), Inches(2.0),
             fill=C_WARN_BG, line_color=RGBColor(0x7f, 0x1d, 0x1d), line_width_pt=0.75)
    add_textbox(s, "概要  |  Applied Materials (AMAT)", Inches(0.65), Inches(1.55),
                Inches(5.2), Inches(0.28), font_size=10, bold=True, color=C_RED)
    add_textbox(s,
                "米国最大の半導体製造装置メーカー。エンティティリスト掲載のSMIC（中国）への\n"
                "装置出荷疑惑でDOJ（米司法省）が捜査。2024年に$153M+の民事和解合意が報道。",
                Inches(0.65), Inches(1.87), Inches(5.2), Inches(1.4),
                font_size=9.5, color=C_GRAY)

    # Stats under case summary
    stat_box(s, Inches(0.5), Inches(3.55), Inches(2.65), Inches(0.95),
             "$153M+", "民事和解合意額\n（2024年報道）", num_color=C_RED, is_dark=True)
    stat_box(s, Inches(3.35), Inches(3.55), Inches(2.65), Inches(0.95),
             "刑事", "法人・個人への\n刑事罰リスク", num_color=C_RED, is_dark=True)

    # Right: timeline + impact
    add_textbox(s, "事件の経緯", Inches(6.3), Inches(1.45), Inches(2.5), Inches(0.25),
                font_size=8, bold=True, color=C_GRAY, font_name="Arial")
    timeline = [
        ("2020-21", "SMICがEntity List掲載。AMAT第三国経由出荷継続の疑惑が発生。"),
        ("2022-23", "DOJ・商務省が調査開始。株主訴訟も並行発生。株価・経営陣に影響。"),
        ("2024",    "$153M+和解合意が報道。コンプライアンス体制の全面見直しを実施。"),
    ]
    ty = Inches(1.75)
    for year, text in timeline:
        add_rect(s, Inches(6.3), ty + Inches(0.05), Inches(0.07), Inches(0.07), fill=C_RED)
        add_textbox(s, year, Inches(6.45), ty, Inches(0.9), Inches(0.25),
                    font_size=9, bold=True, color=C_RED, font_name="Arial")
        add_textbox(s, text, Inches(7.45), ty, Inches(5.3), Inches(0.38),
                    font_size=9.5, color=C_GRAY)
        ty += Inches(0.55)

    impact = [
        ("⚖️", "刑事罰",   "個人・法人への重大罰金。代表者の監督責任。"),
        ("🚫", "業務停止", "輸出業務停止。サプライチェーン即時停止。"),
        ("📉", "信頼毀損", "株価急落・顧客離反・採用困難。"),
    ]
    ix = Inches(6.3)
    for icon, title, body in impact:
        add_rect(s, ix, Inches(3.5), Inches(2.2), Inches(1.05),
                 fill=C_WARN_BG, line_color=RGBColor(0x7f, 0x1d, 0x1d), line_width_pt=0.5)
        add_textbox(s, icon, ix + Inches(0.1), Inches(3.55), Inches(0.3), Inches(0.3), font_size=14, color=C_TEXT_W)
        add_textbox(s, title, ix + Inches(0.42), Inches(3.55), Inches(1.7), Inches(0.25),
                    font_size=10, bold=True, color=C_TEXT_W)
        add_textbox(s, body, ix + Inches(0.1), Inches(3.82), Inches(2.05), Inches(0.65),
                    font_size=8.5, color=C_GRAY)
        ix += Inches(2.35)

    add_footer(s, is_dark=True, slide_num="04 / 15")


def slide_05_asis(prs: Presentation) -> None:
    """As-Is: 現在の体制 — LIGHT"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_LIGHT)

    add_label_title_body(
        s, "AS-IS — 現状分析",
        "典型的な輸出管理体制——潜む5つのリスク",
        "「機能しているように見える」フローの中に、見えていない危機が積み重なっている",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=False,
        accent_color=C_RED,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_RED)

    # Flow diagram
    add_textbox(s, "典型的な現状フロー", Inches(0.5), Inches(1.45), Inches(3), Inches(0.25),
                font_size=8, bold=True, color=C_BLUE, font_name="Arial")
    flow_steps = [
        ("新規取引\n発生",   C_BLUE,                    RGBColor(0xc3, 0xd9, 0xf0)),
        ("手動調査\n2〜5日", C_RED,                     RGBColor(0xfe, 0xca, 0xca)),
        ("Excel記録\nバージョン管理なし", C_RED,         RGBColor(0xfe, 0xca, 0xca)),
        ("外部委託\n2〜8週・50〜200万", C_RED,           RGBColor(0xfe, 0xca, 0xca)),
        ("形骸化承認\n内容理解なし", C_RED,              RGBColor(0xfe, 0xca, 0xca)),
        ("出荷\n← 問題が潜む", RGBColor(0x15, 0x80, 0x3d), RGBColor(0xdc, 0xfc, 0xe7)),
    ]
    fw = Inches(1.85)
    fx = Inches(0.5)
    fy = Inches(1.75)
    for i, (label, text_color, bg_color) in enumerate(flow_steps):
        add_rect(s, fx, fy, fw, Inches(0.7),
                 fill=bg_color, line_color=RGBColor(0xd8, 0xe4, 0xf0), line_width_pt=0.75)
        add_textbox(s, label, fx + Inches(0.05), fy + Inches(0.05), fw - Inches(0.1), Inches(0.62),
                    font_size=9, bold=False, color=text_color, align=PP_ALIGN.CENTER)
        if i < len(flow_steps) - 1:
            add_textbox(s, "→", fx + fw, fy + Inches(0.2), Inches(0.2), Inches(0.3),
                        font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)
        fx += Inches(2.06)

    # 3 risk cards
    risk_cards = [
        ("👤", "人材リスク", "専門知識の属人化",
         "該非判定のエキスパートは国内でも希少。担当者の退職・異動で輸出管理機能が一夜にして消える。後任育成に2〜3年。"),
        ("💸", "コストリスク", "外部委託の高コスト・低速",
         "単発委託で50〜200万円/件、2〜8週間。年間100件なら2〜5億円規模。急ぎ案件に対応不能。機会損失も大きい。"),
        ("🕳️", "規制カバレッジ", "キャッチオール規制の盲点",
         "リスト規制外の汎用品でも輸出令第4条が適用。用途・仕向地・パターンの複合判断は専門家でなければ困難。"),
    ]
    cw = Inches(4.0)
    cx = Inches(0.5)
    cy = Inches(2.6)
    ch = Inches(3.75)
    for icon, tag, title, body in risk_cards:
        card_light(s, cx, cy, cw, ch, icon, tag, title, body,
                   border_color=RGBColor(0xfe, 0xca, 0xca),
                   bg=RGBColor(0xff, 0xf5, 0xf5))
        cx += Inches(4.17)

    add_footer(s, is_dark=False, slide_num="05 / 15")


def slide_06_pain(prs: Presentation) -> None:
    """3つの根本課題 — NAVY"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_NAVY)

    add_label_title_body(
        s, "ROOT CAUSE — 根本課題",
        "3つの根本課題は「連鎖」している",
        "それぞれが単独の問題ではなく、互いを悪化させる負のスパイラル",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=True,
        accent_color=C_RED,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_RED)

    pain_items = [
        ("👤", "課題 1", C_RED,  "専門知識の属人化",
         "外為法＋ECCN＋HSを体系的に理解する人材は国内で極めて希少。担当者=ナレッジ。退職=機能消失。後任育成に平均2〜3年。",
         "↓ だから外部委託に頼る", C_RED),
        ("💸", "課題 2", C_GOLD, "外部委託の限界",
         "年間100件×50万円=5,000万円。2〜8週間の待ち時間。知識が社内に蓄積しない。件数を絞るほどリスクが高まる悪循環。",
         "↓ だから判断を省略する", C_GOLD),
        ("🕳️", "課題 3", C_RED, "キャッチオール規制の見落とし",
         "汎用品でも「用途×仕向地×パターン」の複合判断でキャッチオール発動。専門知識なしの自己判断が見落としを生む。",
         "↓ AMATのような事例につながる", C_RED),
    ]
    cw = Inches(4.0)
    cx = Inches(0.5)
    cy = Inches(1.45)
    card_h = Inches(4.5)
    for icon, tag, tag_color, title, body, note, note_color in pain_items:
        add_rect(s, cx, cy, cw, card_h,
                 fill=C_WARN_BG, line_color=RGBColor(0x7f, 0x1d, 0x1d), line_width_pt=0.75)
        add_textbox(s, icon, cx + Inches(0.15), cy + Inches(0.12), Inches(0.4), Inches(0.4), font_size=24, color=C_TEXT_W)
        add_textbox(s, tag, cx + Inches(0.58), cy + Inches(0.15), cw - Inches(0.7), Inches(0.25),
                    font_size=9, bold=True, color=tag_color, font_name="Arial")
        add_textbox(s, title, cx + Inches(0.15), cy + Inches(0.52), cw - Inches(0.3), Inches(0.45),
                    font_size=14, bold=True, color=C_TEXT_W)
        add_textbox(s, body, cx + Inches(0.15), cy + Inches(1.05), cw - Inches(0.3), Inches(2.0),
                    font_size=10, color=C_GRAY)
        # note box
        add_rect(s, cx, cy + card_h - Inches(0.4), cw, Inches(0.38),
                 fill=RGBColor(0x1a, 0x08, 0x08), line_color=note_color, line_width_pt=0.5)
        add_textbox(s, note, cx + Inches(0.1), cy + card_h - Inches(0.38), cw - Inches(0.2), Inches(0.35),
                    font_size=9, bold=True, color=note_color, align=PP_ALIGN.CENTER)
        cx += Inches(4.17)

    add_footer(s, is_dark=True, slide_num="06 / 15")


def slide_07_tobe(prs: Presentation) -> None:
    """To-Be: 目指すべき姿 — LIGHT"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_LIGHT)

    add_label_title_body(
        s, "TO-BE — VISION",
        "輸出管理を「個人のスキル」から「組織の能力」へ",
        "担当者が変わっても、業務量が増えても、規制が変わっても——常に正しい判断ができる組織",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=False,
        accent_color=RGBColor(0x15, 0x80, 0x3d),
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1),
             fill=RGBColor(0x22, 0xc5, 0x5e))

    # Vision quote box
    add_rect(s, Inches(0.5), Inches(1.45), Inches(12.33), Inches(0.85),
             fill=RGBColor(0xf0, 0xf7, 0xff), line_color=RGBColor(0xc3, 0xd9, 0xf0), line_width_pt=0.75)
    add_textbox(s,
                '"担当者が変わっても、業務量が増えても、規制が変わっても——\n常に正しい判断が、組織のどこからでもできる。"',
                Inches(0.7), Inches(1.52), Inches(11.93), Inches(0.72),
                font_size=13, bold=True, color=C_TEXT_D, align=PP_ALIGN.CENTER)

    # 4 To-Be cards
    tobe_cards = [
        ("🎯", "誰でも使える", "ガイド付きワークフロー",
         "専門家でなくてもAIのガイドに従えば正しい判断に辿り着ける。専門知識に依存しない体制。"),
        ("⚡", "常時稼働", "24/7リアルタイム監視",
         "制裁リスト・規制更新を自動取得。人が寝ている間も24時間対応。ダウンタイムゼロ。"),
        ("📋", "証跡が残る", "全判断の自動記録",
         "根拠・担当者・日時・参照条文を自動保存。監査対応の工数がゼロに。規制当局も納得。"),
        ("📈", "スケールする", "件数増加にコスト比例しない",
         "プラットフォームで100件も1,000件も同じコスト構造。成長しても輸出管理がボトルネックにならない。"),
    ]
    cw = Inches(2.9)
    cx = Inches(0.5)
    cy = Inches(2.4)
    ch = Inches(4.15)
    for icon, tag, title, body in tobe_cards:
        card_light(s, cx, cy, cw, ch, icon, tag, title, body,
                   border_color=RGBColor(0xbb, 0xf7, 0xd0),
                   bg=RGBColor(0xf0, 0xff, 0xf4))
        cx += Inches(3.08)

    add_footer(s, is_dark=False, slide_num="07 / 15")


def slide_08_gap(prs: Presentation) -> None:
    """Gap Analysis — NAVY"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_NAVY)

    add_label_title_body(
        s, "GAP ANALYSIS",
        "なぜ従来手法では「To-Be」に届かないか",
        "ルールベース・専門家採用・外部委託・一般AI——それぞれの限界",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=True,
        accent_color=C_GOLD,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_GOLD)

    # Table header
    headers = ["アプローチ", "速度", "精度", "スケール", "説明可能性", "日本法令"]
    col_widths = [Inches(2.2), Inches(1.9), Inches(2.0), Inches(1.7), Inches(1.9), Inches(1.9)]
    tx = Inches(0.5)
    ty = Inches(1.45)
    hh = Inches(0.32)
    add_rect(s, Inches(0.5), ty, sum(col_widths), hh, fill=RGBColor(0x07, 0x10, 0x20))
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        add_textbox(s, h, tx, ty + Inches(0.05), w, hh - Inches(0.05),
                    font_size=9, bold=True, color=C_GRAY, font_name="Arial")
        tx += w

    rows = [
        ("ルールベースシステム", "◎ 速い",    "△ 規制の組み合わせ爆発に対応不能",
         "◎", "◎ 根拠明示", "✗ 手動更新が追いつかない"),
        ("専門家採用",          "✗ 属人化・育成2〜3年", "◎ 高精度",
         "✗ 人数依存", "◎", "◎"),
        ("外部委託",            "✗ 2〜8週間", "◎",
         "✗ 費用比例増", "◎", "◎"),
        ("一般AIツール",        "◎ 速い",    "△ ハルシネーション",
         "◎", "✗ ブラックボックス", "✗ 日本法令の理解なし"),
    ]
    row_h = Inches(0.65)
    ry = ty + hh
    for row in rows:
        add_rect(s, Inches(0.5), ry, sum(col_widths), row_h,
                 fill=C_CARD_D, line_color=C_BORDER, line_width_pt=0.3)
        rx = Inches(0.5)
        for i, (val, w) in enumerate(zip(row, col_widths)):
            color = C_TEXT_W if i == 0 else (
                C_GREEN if val.startswith("◎") else
                C_RED   if val.startswith("✗") else
                C_GOLD)
            add_textbox(s, val, rx + Inches(0.05), ry + Inches(0.07),
                        w - Inches(0.1), row_h - Inches(0.1),
                        font_size=9, bold=(i == 0), color=color)
            rx += w
        ry += row_h

    # AI Trade Management highlighted row
    add_rect(s, Inches(0.5), ry, sum(col_widths), row_h,
             fill=RGBColor(0x06, 0x10, 0x2a), line_color=C_BLUE, line_width_pt=1.0)
    atm_vals = ["AI Trade Management", "◎ 数秒〜分", "◎ 91%（実測）", "◎ 無制限", "◎ 法的根拠付き", "◎ フル実装"]
    rx = Inches(0.5)
    for i, (val, w) in enumerate(zip(atm_vals, col_widths)):
        color = C_BLUE_LT if i == 0 else C_GREEN
        add_textbox(s, val, rx + Inches(0.05), ry + Inches(0.07),
                    w - Inches(0.1), row_h - Inches(0.1),
                    font_size=9, bold=True, color=color)
        rx += w

    add_footer(s, is_dark=True, slide_num="08 / 15")


def slide_09_value_prop(prs: Presentation) -> None:
    """Value Proposition — DARK CLIMAX"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_DARK)

    # Center glow (simulated with semi-transparent rect)
    add_rect(s, Inches(3.5), Inches(1.5), Inches(6.5), Inches(3.5),
             fill=RGBColor(0x07, 0x0f, 0x22))

    add_textbox(s, "VALUE PROPOSITION",
                Inches(0.5), Inches(0.55), Inches(12.33), Inches(0.3),
                font_size=9, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER, font_name="Arial")

    add_textbox(s, "AIと法の融合で、",
                Inches(0.5), Inches(1.0), Inches(12.33), Inches(0.65),
                font_size=36, bold=True, color=C_TEXT_W, align=PP_ALIGN.CENTER)
    add_textbox(s, "輸出管理の民主化を実現する。",
                Inches(0.5), Inches(1.6), Inches(12.33), Inches(0.65),
                font_size=36, bold=True, color=C_BLUE_LT, align=PP_ALIGN.CENTER)

    add_textbox(s,
                "NeuroSymbolicアーキテクチャ——AIの柔軟な理解力 × 法律規則の確定的判断力 = 「説明できるコンプライアンス判断」",
                Inches(1.0), Inches(2.35), Inches(11.33), Inches(0.4),
                font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 3 value pillars
    pillars = [
        ("Speed",        "8週間 → 数秒",           "外部委託の処理時間を100分の1以下に",
         RGBColor(0x0a, 0x1a, 0x4a), C_BLUE_LT),
        ("Accuracy",     "HS精度91%・2231カテゴリ", "外為法×ECCN×HSの三位一体判定",
         RGBColor(0x07, 0x1c, 0x0e), C_GREEN),
        ("Auditability", "全判断に法的根拠・参照条文", "監査証跡を自動生成。説明責任を果たせる",
         RGBColor(0x1a, 0x14, 0x05), C_GOLD),
    ]
    pw = Inches(3.8)
    px = Inches(0.5)
    py = Inches(3.0)
    ph = Inches(3.25)
    for label, kpi, body, bg, color in pillars:
        add_rect(s, px, py, pw, ph,
                 fill=bg, line_color=color, line_width_pt=1.0)
        add_textbox(s, label, px, py + Inches(0.25), pw, Inches(0.5),
                    font_size=26, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Arial")
        add_textbox(s, kpi, px + Inches(0.1), py + Inches(0.82), pw - Inches(0.2), Inches(0.38),
                    font_size=13, bold=True, color=C_TEXT_W, align=PP_ALIGN.CENTER)
        add_textbox(s, body, px + Inches(0.15), py + Inches(1.28), pw - Inches(0.3), Inches(0.5),
                    font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER)
        px += Inches(4.0)

    add_footer(s, is_dark=True, slide_num="09 / 15")


def slide_10_product(prs: Presentation) -> None:
    """プロダクト全体像 — LIGHT"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_LIGHT)

    add_label_title_body(
        s, "PRODUCT — PLATFORM ARCHITECTURE",
        "エンドツーエンドで輸出管理をカバーする8モジュール",
        "判定 → スクリーニング → 文書 → 監査まで、一貫したワークフローを提供",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=False,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_BLUE)

    # Platform Core bar
    add_rect(s, Inches(0.5), Inches(1.45), Inches(12.33), Inches(0.65),
             fill=RGBColor(0xf0, 0xf7, 0xff), line_color=RGBColor(0xc3, 0xd9, 0xf0), line_width_pt=1.0)
    add_textbox(s, "🏛️", Inches(0.65), Inches(1.52), Inches(0.4), Inches(0.4), font_size=18, color=C_TEXT_D)
    add_textbox(s, "Platform Core", Inches(1.1), Inches(1.52), Inches(2.0), Inches(0.28),
                font_size=11, bold=True, color=C_BLUE)
    add_textbox(s,
                "統合認証・監査ログ・モジュール間連携・全操作の自動証跡保存",
                Inches(1.1), Inches(1.77), Inches(7), Inches(0.25),
                font_size=9, color=RGBColor(0x4a, 0x55, 0x68))
    for tag in ["JWT認証", "Audit Log", "API Gateway"]:
        add_rect(s, Inches(9.5 if tag == "JWT認証" else (10.3 if tag == "Audit Log" else 11.1)),
                 Inches(1.55), Inches(0.75), Inches(0.3),
                 fill=RGBColor(0xe0, 0xee, 0xff), line_color=C_BLUE, line_width_pt=0.5)
        add_textbox(s, tag,
                    Inches(9.5 if tag == "JWT認証" else (10.3 if tag == "Audit Log" else 11.1)),
                    Inches(1.57), Inches(0.75), Inches(0.28),
                    font_size=8, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER, font_name="Arial")

    # 8 modules in 4×2 grid
    modules = [
        ("🔍", "ai_validation",   "外為法該当性判定\n2231カテゴリ対応",           "外為法",  C_BLUE),
        ("🏷️", "ai_classification","HS/ECCN自動分類\n精度91%",                  "HS・ECCN", RGBColor(0x22, 0xc5, 0x5e)),
        ("🚨", "screening",       "OFAC+BIS制裁チェック\n月次自動更新・1秒/件",    "制裁",   C_RED),
        ("🔬", "rnd_assessment",  "R&D技術の輸出規制評価\nNeuroSymbolicリスクスコア","R&D",   C_GOLD),
        ("📑", "patent_search",   "特許出願人制裁チェック\nJ-PlatPat連携",        "特許",   C_BLUE),
        ("📦", "hs_classifier",   "専用HS分類エンジン\nECCN自動付加",             "HS専用", RGBColor(0x22, 0xc5, 0x5e)),
        ("⚖️", "dap",             "デュアルユース判定\nClaude API統合",           "キャッチオール", C_GOLD),
        ("🔗", "API連携",         "ERP・既存システム統合\nCSV/API入力対応",        "Integration", C_GRAY),
    ]
    mw = Inches(2.8)
    mh = Inches(1.3)
    for i, (icon, name, desc, tag, tag_color) in enumerate(modules):
        row = i // 4
        col = i % 4
        mx = Inches(0.5) + col * Inches(3.08)
        my = Inches(2.2) + row * Inches(1.45)
        bg = RGBColor(0xf8, 0xfb, 0xff) if i == 7 else C_WHITE
        add_rect(s, mx, my, mw, mh,
                 fill=bg, line_color=C_LIGHT_BORDER, line_width_pt=0.75)
        add_textbox(s, icon, mx + Inches(0.1), my + Inches(0.1), Inches(0.35), Inches(0.35), font_size=16, color=C_TEXT_D)
        add_textbox(s, name, mx + Inches(0.48), my + Inches(0.12), mw - Inches(0.6), Inches(0.25),
                    font_size=9, bold=True, color=C_BLUE)
        add_textbox(s, desc, mx + Inches(0.1), my + Inches(0.48), mw - Inches(0.2), Inches(0.55),
                    font_size=8.5, color=RGBColor(0x4a, 0x55, 0x68))
        # tag
        add_rect(s, mx + Inches(0.1), my + mh - Inches(0.3), Inches(1.2), Inches(0.24),
                 fill=RGBColor(0xe0, 0xee, 0xff) if tag_color == C_BLUE else RGBColor(0xf0, 0xff, 0xf4),
                 line_color=tag_color, line_width_pt=0.5)
        add_textbox(s, tag, mx + Inches(0.1), my + mh - Inches(0.3), Inches(1.2), Inches(0.24),
                    font_size=7.5, bold=True, color=tag_color, align=PP_ALIGN.CENTER, font_name="Arial")

    add_footer(s, is_dark=False, slide_num="10 / 15")


def slide_11_differentiation(prs: Presentation) -> None:
    """競合優位性 — NAVY"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_NAVY)

    add_label_title_body(
        s, "COMPETITIVE ADVANTAGE",
        "なぜAI Trade Managementが選ばれるか",
        "技術・ドメイン・カバレッジ——3つの同時実現が他にない",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=True,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_BLUE)

    diff_cards = [
        ("🧠", "優位性 1", C_BLUE,
         "NeuroSymbolicアーキテクチャ",
         "AIの柔軟な文書理解 × 法律規則エンジンの確定的判断。説明可能なコンプライアンス。根拠が取れる判断で監査・内部報告に対応。",
         ["FAISS 3層構造", "Symbolic Engine", "Claude API"],
         C_BLUE),
        ("🇯🇵", "優位性 2", C_GREEN,
         "日本法令の深い実装",
         "外為法（輸出令・貨物等省令・外為令）フル実装。HS日本版＋ECCN対照マッピング。経産省更新に連動する継続的メンテナンス。",
         ["外為法フル実装", "2231カテゴリ", "HS×ECCN連携"],
         RGBColor(0x22, 0xc5, 0x5e)),
        ("🔗", "優位性 3", C_GOLD,
         "エンドツーエンドカバレッジ",
         "製品判定 → 取引先スクリーニング → R&D評価 → 特許調査 → 監査ログまで一貫。既存業務フローへのAPI統合も対応。",
         ["8モジュール統合", "API統合対応", "監査証跡自動化"],
         C_GOLD),
    ]
    cw = Inches(4.0)
    cx = Inches(0.5)
    cy = Inches(1.45)
    ch = Inches(5.3)
    for icon, num_tag, tag_color, title, body, tags, accent in diff_cards:
        add_rect(s, cx, cy, cw, ch,
                 fill=C_GOOD_BG, line_color=accent, line_width_pt=1.0)
        add_textbox(s, icon, cx + Inches(0.15), cy + Inches(0.12), Inches(0.5), Inches(0.5), font_size=26, color=C_TEXT_W)
        add_textbox(s, num_tag, cx + Inches(0.68), cy + Inches(0.18), cw - Inches(0.8), Inches(0.25),
                    font_size=9, bold=True, color=tag_color, font_name="Arial")
        add_textbox(s, title, cx + Inches(0.15), cy + Inches(0.62), cw - Inches(0.3), Inches(0.55),
                    font_size=14, bold=True, color=C_TEXT_W)
        add_textbox(s, body, cx + Inches(0.15), cy + Inches(1.25), cw - Inches(0.3), Inches(2.2),
                    font_size=10, color=C_GRAY)
        # tags
        ty = cy + ch - Inches(0.5)
        tx = cx + Inches(0.15)
        for tag in tags:
            tw = Inches(1.3)
            add_rect(s, tx, ty, tw, Inches(0.3),
                     fill=RGBColor(0x06, 0x10, 0x28), line_color=accent, line_width_pt=0.5)
            add_textbox(s, tag, tx, ty, tw, Inches(0.3),
                        font_size=8, bold=True, color=accent, align=PP_ALIGN.CENTER, font_name="Arial")
            tx += Inches(1.38)
        cx += Inches(4.17)

    add_footer(s, is_dark=True, slide_num="11 / 15")


def slide_12_proof(prs: Presentation) -> None:
    """実証データ — LIGHT"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_LIGHT)

    add_label_title_body(
        s, "EVIDENCE — 実証データ",
        "主観ではなく数字で語る——実装済みシステムの実測値",
        "M4 Mac最適化・FAISS再構築・Layer C精度改善を含む最新実績",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=False,
        accent_color=RGBColor(0x15, 0x80, 0x3d),
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1),
             fill=RGBColor(0x22, 0xc5, 0x5e))

    # 4 stats
    proof_stats = [
        ("91%",   "HS判定精度\n（Layer C最適化後）",    C_GREEN),
        ("2,231", "外為法分類カテゴリ\n完全対応数",      C_BLUE),
        ("<1秒",  "制裁スクリーニング速度\nOFAC+BIS同時照合", C_GREEN),
        ("100×",  "処理速度改善\n人手比較（8週→数秒）",  C_BLUE),
    ]
    sw = Inches(2.9)
    sx = Inches(0.5)
    for number, label, color in proof_stats:
        stat_box(s, sx, Inches(1.45), sw, Inches(1.1), number, label, num_color=color, is_dark=False)
        sx += Inches(3.05)

    # Left: progress bars
    add_textbox(s, "精度改善の軌跡", Inches(0.5), Inches(2.7), Inches(5), Inches(0.25),
                font_size=9, bold=True, color=C_BLUE, font_name="Arial")
    add_rect(s, Inches(0.5), Inches(2.95), Inches(5.5), Inches(2.3),
             fill=C_WHITE, line_color=C_LIGHT_BORDER, line_width_pt=0.75)

    for label, pct, color, bar_bg in [
        ("Layer C (HS) 改善前: 83%", 0.83, C_RED, RGBColor(0xfe, 0xe2, 0xe2)),
        ("Layer C (HS) 改善後: 91%", 0.91, RGBColor(0x22, 0xc5, 0x5e), RGBColor(0xdc, 0xfc, 0xe7)),
    ]:
        py_off = Inches(0.3) if "改善前" in label else Inches(1.1)
        add_textbox(s, label, Inches(0.65), Inches(2.95) + py_off, Inches(5.2), Inches(0.28),
                    font_size=10, bold=True, color=C_TEXT_D)
        add_rect(s, Inches(0.65), Inches(2.95) + py_off + Inches(0.28), Inches(5.0), Inches(0.2),
                 fill=bar_bg, line_color=color, line_width_pt=0.3)
        add_rect(s, Inches(0.65), Inches(2.95) + py_off + Inches(0.28),
                 Inches(5.0 * pct), Inches(0.2), fill=color)

    add_textbox(s, "※ FEFTA汚染ラベル除去・FAISS再構築による根本解決",
                Inches(0.65), Inches(5.1), Inches(5.2), Inches(0.25),
                font_size=8, color=C_GRAY)

    # Right: auto features
    add_textbox(s, "自動化・継続稼働", Inches(6.3), Inches(2.7), Inches(6.5), Inches(0.25),
                font_size=9, bold=True, color=C_BLUE, font_name="Arial")
    auto_items = [
        ("🔄", "制裁リスト月次自動更新", "OFAC/BIS API連携。手動作業ゼロ。",        "自動", C_GREEN),
        ("🌐", "多言語対応 (日英中)",   "multilingual-e5-largeモデル採用。",      "多言語", C_BLUE),
        ("📊", "ECCN付加率95%超",      "HSコードから自動推定。マッピングDB連携。", "高精度", C_GREEN),
        ("🔐", "全操作の監査ログ",     "全判断に担当者・日時・根拠を自動記録。",   "証跡", C_GOLD),
    ]
    aiy = Inches(2.95)
    for icon, title, desc, badge, badge_color in auto_items:
        add_rect(s, Inches(6.3), aiy, Inches(6.5), Inches(0.55),
                 fill=C_WHITE, line_color=C_LIGHT_BORDER, line_width_pt=0.75)
        add_textbox(s, icon, Inches(6.45), aiy + Inches(0.1), Inches(0.35), Inches(0.35), font_size=14, color=C_TEXT_D)
        add_textbox(s, title, Inches(6.85), aiy + Inches(0.08), Inches(4.8), Inches(0.25),
                    font_size=10, bold=True, color=C_TEXT_D)
        add_textbox(s, desc, Inches(6.85), aiy + Inches(0.32), Inches(4.8), Inches(0.2),
                    font_size=8.5, color=RGBColor(0x4a, 0x55, 0x68))
        add_rect(s, Inches(12.0), aiy + Inches(0.12), Inches(0.7), Inches(0.28),
                 fill=RGBColor(0xf0, 0xff, 0xf4) if badge_color == C_GREEN else RGBColor(0xe0, 0xee, 0xff),
                 line_color=badge_color, line_width_pt=0.5)
        add_textbox(s, badge, Inches(12.0), aiy + Inches(0.12), Inches(0.7), Inches(0.28),
                    font_size=7.5, bold=True, color=badge_color, align=PP_ALIGN.CENTER, font_name="Arial")
        aiy += Inches(0.6)

    add_footer(s, is_dark=False, slide_num="12 / 15")


def slide_13_roadmap(prs: Presentation) -> None:
    """導入ロードマップ — LIGHT"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_LIGHT)

    add_label_title_body(
        s, "IMPLEMENTATION ROADMAP",
        "まず30日間の無償パイロット——リスクゼロのスタート",
        "段階的導入で価値を確認しながら拡大。強制・リスクなしのアプローチ。",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=False,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_BLUE)

    phases = [
        ("PHASE 1", "パイロット検証", RGBColor(0x0a, 0x6c, 0xff),
         "特定製品カテゴリ（50〜100品）でHS/ECCN判定をライブ検証。既存判断結果との照合・精度確認。制裁スクリーニングAPIのテスト統合。",
         "Month 1〜2", ["無償", "30日間"]),
        ("PHASE 2", "本格展開", RGBColor(0x22, 0xc5, 0x5e),
         "全製品カテゴリへの拡張。審査フローへの組み込み（承認ワークフロー連携）。担当者トレーニング・運用手順の整備。",
         "Month 3〜4", ["全製品展開"]),
        ("PHASE 3", "システム統合", C_GOLD,
         "ERP・在庫システムとのAPI連携。監査レポート自動生成の設定。効果測定・ROI確定・継続契約の検討。",
         "Month 5〜6", ["ERP連携"]),
    ]
    pw = Inches(3.8)
    px = Inches(0.5)
    ph = Inches(4.2)
    py = Inches(1.45)
    arrows_x = [Inches(0.5) + pw, Inches(0.5) + pw * 2 + Inches(0.37)]
    for i, (phase_num, title, color, body, duration, tags) in enumerate(phases):
        add_rect(s, px, py, pw, ph,
                 fill=C_WHITE, line_color=color, line_width_pt=1.5)
        # Top color bar
        add_rect(s, px, py, pw, Inches(0.06), fill=color)
        add_textbox(s, phase_num, px + Inches(0.15), py + Inches(0.12), pw - Inches(0.3), Inches(0.25),
                    font_size=9, bold=True, color=color, font_name="Arial")
        add_textbox(s, title, px + Inches(0.15), py + Inches(0.38), pw - Inches(0.3), Inches(0.38),
                    font_size=14, bold=True, color=C_TEXT_D)
        for j, tag in enumerate(tags):
            tw = Inches(0.9)
            tx = px + Inches(0.15) + j * Inches(1.0)
            add_rect(s, tx, py + Inches(0.8), tw, Inches(0.25),
                     fill=RGBColor(0xe0, 0xee, 0xff), line_color=color, line_width_pt=0.5)
            add_textbox(s, tag, tx, py + Inches(0.8), tw, Inches(0.25),
                        font_size=7.5, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Arial")
        add_textbox(s, body, px + Inches(0.15), py + Inches(1.12), pw - Inches(0.3), Inches(2.0),
                    font_size=9.5, color=RGBColor(0x4a, 0x55, 0x68))
        add_textbox(s, duration, px + Inches(0.15), py + ph - Inches(0.4), pw - Inches(0.3), Inches(0.3),
                    font_size=10, bold=True, color=C_GOLD)
        # Arrow
        if i < 2:
            add_textbox(s, "→", px + pw, py + Inches(1.8), Inches(0.37), Inches(0.4),
                        font_size=20, color=color, align=PP_ALIGN.CENTER)
        px += Inches(4.17)

    # KPI row
    add_rect(s, Inches(0.5), Inches(5.75), Inches(12.33), Inches(0.55),
             fill=RGBColor(0xf0, 0xf7, 0xff), line_color=RGBColor(0xc3, 0xd9, 0xf0), line_width_pt=0.75)
    add_textbox(s, "KPI目標", Inches(0.65), Inches(5.83), Inches(1.0), Inches(0.25),
                font_size=8, bold=True, color=C_BLUE, font_name="Arial")
    kpis = [
        ("判定時間:", "8週間 → 1日以内", C_BLUE),
        ("外部委託費:", "-70%以上", RGBColor(0x15, 0x80, 0x3d)),
        ("監査証跡:", "100%カバー（全判断記録）", C_BLUE),
    ]
    kx = Inches(1.8)
    for label, val, color in kpis:
        add_textbox(s, f"{label} {val}", kx, Inches(5.84), Inches(3.5), Inches(0.28),
                    font_size=10, bold=False, color=C_TEXT_D)
        add_textbox(s, val, kx + Inches(1.05), Inches(5.84), Inches(2.4), Inches(0.28),
                    font_size=10, bold=True, color=color)
        kx += Inches(4.0)

    add_footer(s, is_dark=False, slide_num="13 / 15")


def slide_14_roi(prs: Presentation) -> None:
    """ROI試算 — LIGHT"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_LIGHT)

    add_label_title_body(
        s, "INVESTMENT & RETURN",
        "コスト削減 + リスク回避価値で考えるROI",
        "中規模製造業（年間輸出取引200件）想定のシミュレーション",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=False,
        accent_color=C_GOLD,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_GOLD)

    # Before column
    add_rect(s, Inches(0.5), Inches(1.45), Inches(5.7), Inches(4.0),
             fill=RGBColor(0xff, 0xf5, 0xf5), line_color=RGBColor(0xfe, 0xca, 0xca), line_width_pt=1.0)
    add_textbox(s, "BEFORE — 現状コスト（年間）", Inches(0.65), Inches(1.55),
                Inches(5.4), Inches(0.28), font_size=9, bold=True, color=C_RED, font_name="Arial")
    before_items = [
        ("外部委託 200件×50万円",  "1億円"),
        ("担当者工数 3名分",        "2,400万円"),
        ("遅延による機会損失（推定）","5,000万円"),
    ]
    iy = Inches(1.9)
    for label, val in before_items:
        add_rect(s, Inches(0.65), iy, Inches(5.4), Pt(0.5), fill=RGBColor(0xfe, 0xca, 0xca))
        add_textbox(s, label, Inches(0.65), iy + Inches(0.06), Inches(3.5), Inches(0.35),
                    font_size=10, color=C_TEXT_D)
        add_textbox(s, val, Inches(4.15), iy + Inches(0.06), Inches(1.9), Inches(0.35),
                    font_size=12, bold=True, color=C_RED, align=PP_ALIGN.RIGHT)
        iy += Inches(0.65)
    add_rect(s, Inches(0.65), iy, Inches(5.4), Pt(1.0), fill=C_RED)
    add_textbox(s, "合計", Inches(0.65), iy + Inches(0.06), Inches(3), Inches(0.38),
                font_size=12, bold=True, color=C_TEXT_D)
    add_textbox(s, "約1.9億円/年", Inches(3.15), iy + Inches(0.06), Inches(2.9), Inches(0.38),
                font_size=16, bold=True, color=C_RED, align=PP_ALIGN.RIGHT)

    # VS
    add_textbox(s, "vs.", Inches(6.3), Inches(3.1), Inches(0.7), Inches(0.5),
                font_size=20, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    # After column
    add_rect(s, Inches(7.1), Inches(1.45), Inches(5.7), Inches(4.0),
             fill=RGBColor(0xf0, 0xff, 0xf4), line_color=RGBColor(0xbb, 0xf7, 0xd0), line_width_pt=1.0)
    add_textbox(s, "AFTER — 導入後コスト（年間）", Inches(7.25), Inches(1.55),
                Inches(5.4), Inches(0.28), font_size=9, bold=True,
                color=RGBColor(0x15, 0x80, 0x3d), font_name="Arial")
    after_items = [
        ("プラットフォーム利用料",     "要相談"),
        ("担当者工数 1名（AI補助）", "800万円"),
        ("遅延機会損失",              "ほぼゼロ"),
    ]
    ay = Inches(1.9)
    for label, val in after_items:
        add_rect(s, Inches(7.25), ay, Inches(5.4), Pt(0.5), fill=RGBColor(0xbb, 0xf7, 0xd0))
        add_textbox(s, label, Inches(7.25), ay + Inches(0.06), Inches(3.5), Inches(0.35),
                    font_size=10, color=C_TEXT_D)
        add_textbox(s, val, Inches(10.75), ay + Inches(0.06), Inches(1.9), Inches(0.35),
                    font_size=12, bold=True, color=RGBColor(0x15, 0x80, 0x3d), align=PP_ALIGN.RIGHT)
        ay += Inches(0.65)
    add_rect(s, Inches(7.25), ay, Inches(5.4), Pt(1.0), fill=RGBColor(0x22, 0xc5, 0x5e))
    add_textbox(s, "合計", Inches(7.25), ay + Inches(0.06), Inches(3), Inches(0.38),
                font_size=12, bold=True, color=C_TEXT_D)
    add_textbox(s, "大幅削減", Inches(10.25), ay + Inches(0.06), Inches(2.9), Inches(0.38),
                font_size=16, bold=True, color=RGBColor(0x15, 0x80, 0x3d), align=PP_ALIGN.RIGHT)

    # Risk avoidance
    add_rect(s, Inches(0.5), Inches(5.55), Inches(12.33), Inches(0.8),
             fill=RGBColor(0xff, 0xfb, 0xeb), line_color=RGBColor(0xfd, 0xba, 0x74), line_width_pt=0.75)
    add_textbox(s, "⚠️", Inches(0.65), Inches(5.62), Inches(0.4), Inches(0.4), font_size=20, color=C_TEXT_D)
    add_textbox(s, "最大のROI — リスク回避価値",
                Inches(1.1), Inches(5.6), Inches(8), Inches(0.28),
                font_size=9, bold=True, color=RGBColor(0x92, 0x40, 0x0e), font_name="Arial")
    add_textbox(s,
                "制裁1件を防げれば、プラットフォームへの投資は瞬時に回収。AMATベースで$153M+（約230億円）のリスク回避価値。",
                Inches(1.1), Inches(5.86), Inches(10.5), Inches(0.3),
                font_size=9.5, color=RGBColor(0x37, 0x41, 0x51))
    add_textbox(s, "∞", Inches(12.1), Inches(5.58), Inches(0.6), Inches(0.5),
                font_size=24, bold=True, color=RGBColor(0x92, 0x40, 0x0e), align=PP_ALIGN.CENTER)

    add_footer(s, is_dark=False, slide_num="14 / 15")


def slide_15_next_steps(prs: Presentation) -> None:
    """次のステップ — DARK CTA"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG_DARK)

    add_label_title_body(
        s, "NEXT STEPS — ACTION",
        "本日30分で輪郭が見えました。次回60分でリアルを体験しましょう。",
        "2週間以内のデモ設定でリスクなく価値を確認できます",
        Inches(0.5), Inches(0.3), Inches(12.33), is_dark=True,
    )
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Pt(1), fill=C_BLUE)

    # Left: next meeting
    add_textbox(s, "次回（2nd Meeting）でお見せすること",
                Inches(0.5), Inches(1.45), Inches(6.5), Inches(0.25),
                font_size=9, bold=True, color=C_BLUE, font_name="Arial")
    next_items = [
        ("01", "ライブデモ",              "御社の実製品データでリアルタイム判定"),
        ("02", "業務フィット&ギャップ分析", "御社の現行フローとの詳細照合"),
        ("03", "カスタムROI試算",          "御社規模に合わせた投資対効果"),
        ("04", "導入条件・技術要件",        "パイロット開始までのステップ確認"),
    ]
    ny = Inches(1.75)
    for num, main, sub in next_items:
        add_rect(s, Inches(0.5), ny, Inches(6.3), Inches(0.6),
                 fill=C_CARD_D, line_color=C_BORDER, line_width_pt=0.75)
        add_textbox(s, num, Inches(0.65), ny + Inches(0.1), Inches(0.4), Inches(0.4),
                    font_size=14, bold=True, color=C_BLUE, font_name="Arial")
        add_textbox(s, main, Inches(1.15), ny + Inches(0.08), Inches(5.5), Inches(0.25),
                    font_size=11, bold=True, color=C_TEXT_W)
        add_textbox(s, sub, Inches(1.15), ny + Inches(0.33), Inches(5.5), Inches(0.22),
                    font_size=9, color=C_GRAY)
        ny += Inches(0.65)

    # Right: CTA
    add_textbox(s, "本日お願いする3つのアクション",
                Inches(7.1), Inches(1.45), Inches(6.0), Inches(0.25),
                font_size=9, bold=True, color=C_BLUE, font_name="Arial")

    add_rect(s, Inches(7.1), Inches(1.75), Inches(6.0), Inches(2.4),
             fill=RGBColor(0x06, 0x10, 0x28), line_color=C_BLUE, line_width_pt=0.75)
    cta_items = [
        "デモ日程の調整（2週間以内で1時間）",
        "業務ヒアリングシートのご回答（5分）",
        "パイロット対象製品カテゴリのご検討",
    ]
    cy = Inches(2.0)
    for item in cta_items:
        add_textbox(s, "✓", Inches(7.3), cy, Inches(0.35), Inches(0.35),
                    font_size=12, bold=True, color=C_GREEN)
        add_textbox(s, item, Inches(7.7), cy, Inches(6.2), Inches(0.32),
                    font_size=10, color=C_TEXT_W)
        cy += Inches(0.62)

    # Free pilot box
    add_rect(s, Inches(7.1), Inches(4.3), Inches(6.0), Inches(2.3),
             fill=RGBColor(0x06, 0x14, 0x0e), line_color=C_GREEN, line_width_pt=1.0)
    add_textbox(s, "無償パイロット 30日間", Inches(7.1), Inches(4.5), Inches(6.0), Inches(0.3),
                font_size=10, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER, font_name="Arial")
    add_textbox(s, "リスクゼロで始められます", Inches(7.1), Inches(4.85), Inches(6.0), Inches(0.5),
                font_size=22, bold=True, color=C_TEXT_W, align=PP_ALIGN.CENTER)
    add_textbox(s, "価値を確認してから判断。コミットメントは一切不要。",
                Inches(7.1), Inches(5.42), Inches(6.0), Inches(0.28),
                font_size=9.5, color=C_GRAY, align=PP_ALIGN.CENTER)

    # Thank you
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.33), Pt(1), fill=C_BORDER)
    add_textbox(s, "AI TRADE COMPLIANCE MANAGEMENT",
                Inches(0.5), Inches(6.65), Inches(8), Inches(0.25),
                font_size=7, bold=True, color=RGBColor(0x2d, 0x37, 0x48), font_name="Arial")
    add_textbox(s, "THANK YOU", Inches(9.8), Inches(6.63), Inches(3), Inches(0.3),
                font_size=11, bold=True, color=C_BLUE, align=PP_ALIGN.RIGHT, font_name="Arial")


# ════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════

def build(output_path: Path) -> None:
    print("🔨  Building sales deck …")
    prs = new_prs()

    builders = [
        slide_01_cover,
        slide_02_agenda,
        slide_03_why_now,
        slide_04_amat,
        slide_05_asis,
        slide_06_pain,
        slide_07_tobe,
        slide_08_gap,
        slide_09_value_prop,
        slide_10_product,
        slide_11_differentiation,
        slide_12_proof,
        slide_13_roadmap,
        slide_14_roi,
        slide_15_next_steps,
    ]
    for i, fn in enumerate(builders, 1):
        print(f"  [{i:02d}/{len(builders)}] {fn.__name__}")
        fn(prs)

    prs.save(str(output_path))
    size_mb = output_path.stat().st_size / 1_048_576
    print(f"\n✅  Saved → {output_path}  ({size_mb:.1f} MB)")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Build AI Trade Management sales deck (.pptx)")
    parser.add_argument("--output", type=Path,
                        default=OUT_DIR / "sales_deck.pptx",
                        help="Output path for the PPTX file")
    args = parser.parse_args()
    build(args.output)
