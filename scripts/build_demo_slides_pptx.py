"""
scripts/build_demo_slides_pptx.py
──────────────────────────────────
デモシナリオ PPTX を生成し Google Drive にアップロード → Google Slides に変換する。

使い方:
    python3 scripts/build_demo_slides_pptx.py

出力:
    data/demo/demo_scenario_slides.pptx  — ローカル PPTX ファイル
    Google Drive に Google Slides として保存（コンソールに URL を出力）
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

# ─── python-pptx ──────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    sys.exit("python-pptx が必要です: pip install python-pptx")

# ─── Google API ───────────────────────────────────────────────────────────────
try:
    from google.oauth2.credentials import Credentials
    from google.auth.transport.requests import Request
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload
except ImportError:
    sys.exit("google-api-python-client が必要です: pip install google-api-python-client")

# ─── Paths ────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "data" / "demo" / "demo_scenario_slides.pptx"
TOKEN_PATH = ROOT / "token.json"
CREDS_PATH = ROOT / "credentials.json"

# ─── Color Palette ────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0f, 0x27, 0x44)   # deep navy
C_BG2      = RGBColor(0x1e, 0x3a, 0x5f)   # navy mid
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_WHITE70  = RGBColor(0xCC, 0xD6, 0xE8)
C_S1       = RGBColor(0xF5, 0x9E, 0x0B)   # amber  — シナリオ1
C_S1_LIGHT = RGBColor(0xFD, 0xE6, 0x8A)
C_S2       = RGBColor(0x25, 0x63, 0xEB)   # blue   — シナリオ2
C_S2_LIGHT = RGBColor(0xBF, 0xDB, 0xFE)
C_S3       = RGBColor(0x16, 0xA3, 0x4A)   # green  — シナリオ3
C_S3_LIGHT = RGBColor(0xBB, 0xF7, 0xD0)
C_RED      = RGBColor(0xDC, 0x26, 0x26)
C_RED_LT   = RGBColor(0xFC, 0xA5, 0xA5)
C_WARN     = RGBColor(0xFD, 0xE6, 0x8A)
C_GRAY     = RGBColor(0x64, 0x74, 0x8B)


# ─── Helpers ──────────────────────────────────────────────────────────────────

def add_bg(slide, color: RGBColor = C_BG):
    """スライド背景を単色で塗りつぶす。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor, *, alpha=None):
    """塗りつぶし矩形を追加。"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_textbox(slide, left, top, width, height, text: str, *,
                font_size=18, bold=False, color: RGBColor = C_WHITE,
                align=PP_ALIGN.LEFT, wrap=True) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_multiline(slide, left, top, width, height, lines: list[tuple], *,
                  spacing_pt: float = 4):
    """
    lines = [(text, font_size, bold, color, align), ...]
    複数行を1テキストボックスに収める。
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, fsize, fbold, fcolor, falign) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = falign
        p.space_after = Pt(spacing_pt)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fsize)
        run.font.bold = fbold
        run.font.color.rgb = fcolor


# ─── Slide builders ───────────────────────────────────────────────────────────

def build_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, C_BG)

    # Accent stripe (top)
    add_rect(slide, 0, 0, 10, 0.06, C_S2)

    # Subtitle small
    add_textbox(slide, 0.5, 0.5, 9, 0.4,
                "AI Trade Management Platform  —  Demo Scenario",
                font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    # Main title
    add_multiline(slide, 0.5, 1.0, 9, 1.6, [
        ("輸出管理 AI デモシナリオ解説", 36, True, C_WHITE, PP_ALIGN.CENTER),
        ("実業務シナリオで体感する — 審査フローと AI 判定の連携", 14, False, C_WHITE70, PP_ALIGN.CENTER),
    ])

    # 3 scenario pills as colored boxes
    pill_data = [
        (0.5,  3.0, C_S1, C_S1_LIGHT, "SCENARIO 01", "R&D 対中リスク起案",   "rnd_assessment  :8003"),
        (3.8,  3.0, C_S2, C_S2_LIGHT, "SCENARIO 02", "CMP スラリー HS 付番", "ai_classification :8002"),
        (7.1, 3.0, C_S3, C_S3_LIGHT, "SCENARIO 03", "半導体メモリ キャッチオール審査", "ai_validation :8001"),
    ]
    for x, y, c_main, c_light, num, title, module in pill_data:
        # pill background
        shape = slide.shapes.add_shape(1,
            Inches(x), Inches(y), Inches(2.6), Inches(1.5))
        shape.line.color.rgb = c_main
        shape.line.width = Pt(1.5)
        shape.fill.solid()
        # RGBColor stores values as a 6-digit hex string; parse manually
        r = int(str(c_main)[0:2], 16)
        g = int(str(c_main)[2:4], 16)
        b = int(str(c_main)[4:6], 16)
        shape.fill.fore_color.rgb = RGBColor(
            min(r + 20, 255),
            min(g + 15, 255),
            min(b + 30, 255),
        )
        # text inside pill
        add_multiline(slide, x + 0.1, y + 0.1, 2.4, 1.3, [
            (num,    9,  True,  c_main,  PP_ALIGN.LEFT),
            (title,  11, True,  c_light, PP_ALIGN.LEFT),
            (module,  8, False, C_GRAY,  PP_ALIGN.LEFT),
        ])

    # Bottom meta
    add_textbox(slide, 0.5, 5.1, 9, 0.3,
                "2026-03-21  |  AI_TradeManagement v1.0  |  For Demo Use Only",
                font_size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

    # Bottom stripe
    add_rect(slide, 0, 5.55, 10, 0.07, C_BG2)


def build_s1_slide(prs: Presentation):
    """シナリオ1: R&D 対中国リスク起案"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x12, 0x0a, 0x02))

    # Header bar
    add_rect(slide, 0, 0, 10, 0.85, RGBColor(0x1a, 0x10, 0x04))
    add_rect(slide, 0, 0.85, 10, 0.03, C_S1)

    add_multiline(slide, 0.3, 0.08, 7.5, 0.7, [
        ("SCENARIO  01", 9, True, C_S1, PP_ALIGN.LEFT),
        ("R&D 案件起案 — 対中国リスク半導体共同研究", 18, True, C_S1_LIGHT, PP_ALIGN.LEFT),
    ])
    add_textbox(slide, 7.8, 0.2, 2.0, 0.45,
                "rnd_assessment\n:8003",
                font_size=9, bold=True, color=C_S1, align=PP_ALIGN.RIGHT)

    # ── Left column ──
    col1_x = 0.3

    # Box 1: プロジェクト概要
    add_rect(slide, col1_x, 1.0, 4.5, 2.0, RGBColor(0x28, 0x18, 0x06))
    add_textbox(slide, col1_x + 0.1, 1.05, 4.3, 0.25,
                "▌ プロジェクト概要", font_size=9, bold=True, color=C_S1)
    add_multiline(slide, col1_x + 0.15, 1.35, 4.2, 1.55, [
        ("案件名: 次世代3nm FinFET プロセス開発",          10, False, C_WHITE70, PP_ALIGN.LEFT),
        ("相手方: 深圳半導体先端研究院（SSARI）",           10, False, C_WHITE70, PP_ALIGN.LEFT),
        ("開示技術: CMPレシピ・装置制御アルゴリズム",        10, False, C_WHITE70, PP_ALIGN.LEFT),
        ("対象ECCN: 3B001（半導体製造装置）",               10, False, C_S1_LIGHT, PP_ALIGN.LEFT),
        ("外為法: EL-3（電子部品・集積回路）",               10, False, C_S1_LIGHT, PP_ALIGN.LEFT),
    ])

    # Box 2: みなし輸出人物
    add_rect(slide, col1_x, 3.1, 4.5, 1.65, RGBColor(0x28, 0x18, 0x06))
    add_textbox(slide, col1_x + 0.1, 3.15, 4.3, 0.25,
                "▌ みなし輸出 登録人物", font_size=9, bold=True, color=C_S1)
    add_multiline(slide, col1_x + 0.15, 3.43, 4.2, 1.2, [
        ("王 建国（WANG Jianguo）/ SSARI 材料科学部",   10, True,  C_WHITE,   PP_ALIGN.LEFT),
        ("国籍: CN  ｜  二重雇用: COSTIND系機関",        9,  False, C_RED_LT,  PP_ALIGN.LEFT),
        ("ECCN 3B001 アクセス  ｜  カテゴリ A（要許可）", 9,  False, C_RED_LT,  PP_ALIGN.LEFT),
        ("日本在留年数: 0年  ｜  MCF政策下機関所属",       9,  False, C_WARN,    PP_ALIGN.LEFT),
    ])

    # ── Right column ──
    col2_x = 5.1

    # Box 3: 操作フロー
    add_rect(slide, col2_x, 1.0, 4.5, 1.2, RGBColor(0x20, 0x14, 0x05))
    add_textbox(slide, col2_x + 0.1, 1.05, 4.3, 0.25,
                "▌ 操作フロー", font_size=9, bold=True, color=C_S1)
    add_multiline(slide, col2_x + 0.15, 1.33, 4.2, 0.8, [
        ("① 案件情報入力  →  ② AI リスク評価実行", 10, False, C_WHITE70, PP_ALIGN.LEFT),
        ("③ みなし輸出人物登録  →  ④ 判定結果確認", 10, False, C_WHITE70, PP_ALIGN.LEFT),
    ])

    # Box 4: AI評価結果
    add_rect(slide, col2_x, 2.3, 4.5, 1.6, RGBColor(0x20, 0x14, 0x05))
    add_textbox(slide, col2_x + 0.1, 2.35, 4.3, 0.25,
                "▌ AI 評価結果（期待値）", font_size=9, bold=True, color=C_S1)
    add_multiline(slide, col2_x + 0.15, 2.63, 4.2, 1.2, [
        ("技術リスク: ██ HIGH（3nm FinFET = 先端プロセス）",   10, False, C_RED_LT, PP_ALIGN.LEFT),
        ("需要者リスク: ██ HIGH（MCF機関・軍民融合）",          10, False, C_RED_LT, PP_ALIGN.LEFT),
        ("規制リスク: ██ HIGH（CN = D:1/D:3/D:4/D:5）",       10, False, C_RED_LT, PP_ALIGN.LEFT),
    ])

    # Final verdict box
    add_rect(slide, col2_x, 4.0, 4.5, 0.65, RGBColor(0x3b, 0x0c, 0x0c))
    add_textbox(slide, col2_x + 0.15, 4.1, 4.2, 0.45,
                "総合判定:  HIGH RISK — 経済産業大臣許可申請 推奨",
                font_size=11, bold=True, color=C_RED_LT, align=PP_ALIGN.CENTER)

    # DAP alert at bottom
    add_rect(slide, col1_x, 4.75, 9.3, 0.5, RGBColor(0x1a, 0x1a, 0x08))
    add_textbox(slide, col1_x + 0.1, 4.82, 9.0, 0.35,
                "👔 DAP 先輩担当者: 「スクリーニング未実施です。品目管理連携前に取引先照合を実施してください (http://localhost:8005)」",
                font_size=9, color=C_WARN)

    # Footer
    add_rect(slide, 0, 5.5, 10, 0.12, C_S1)
    add_textbox(slide, 0.3, 5.5, 9.4, 0.12,
                "AI_TradeManagement  —  Demo Scenario 01  |  http://localhost:8003",
                font_size=7, color=C_BG)


def build_s2_slide(prs: Presentation):
    """シナリオ2: CMPスラリー HS付番 & EAR判定"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x02, 0x08, 0x18))

    # Header bar
    add_rect(slide, 0, 0, 10, 0.85, RGBColor(0x06, 0x12, 0x28))
    add_rect(slide, 0, 0.85, 10, 0.03, C_S2)

    add_multiline(slide, 0.3, 0.08, 7.5, 0.7, [
        ("SCENARIO  02", 9, True, C_S2, PP_ALIGN.LEFT),
        ("品目登録 — CMPスラリーのHSコード付番 & EAR判定", 18, True, C_S2_LIGHT, PP_ALIGN.LEFT),
    ])
    add_textbox(slide, 7.3, 0.2, 2.5, 0.45,
                "ai_classification\n:8002",
                font_size=9, bold=True, color=C_S2, align=PP_ALIGN.RIGHT)

    # ── Left: BOM Table ──
    col1_x = 0.3

    # Product info
    add_rect(slide, col1_x, 1.0, 4.6, 0.6, RGBColor(0x06, 0x14, 0x30))
    add_textbox(slide, col1_x + 0.1, 1.05, 4.4, 0.2,
                "▌ 品目情報", font_size=9, bold=True, color=C_S2)
    add_textbox(slide, col1_x + 0.15, 1.27, 4.3, 0.3,
                "CMPスラリー（3nm向け シリカ系）  |  HS 3824.99  |  ECCN 1C350（推定）",
                font_size=9.5, color=C_S2_LIGHT)

    # BOM Table header
    add_rect(slide, col1_x, 1.7, 4.6, 0.3, RGBColor(0x1a, 0x3a, 0x6b))
    add_multiline(slide, col1_x + 0.05, 1.72, 4.5, 0.26, [
        ("  原料名              原産国   含有量   ECCN     注意",
         8.5, True, C_S2_LIGHT, PP_ALIGN.LEFT),
    ])

    # BOM rows
    bom_rows = [
        ("  コロイダルシリカ（Cabot）", "US",   "25.0%", "EAR99",  "⚠ de minimis超過",
         RGBColor(0x3b, 0x28, 0x05), C_WARN),
        ("  過酸化水素30%（三菱ガス）", "JP",   " 5.0%", "1C350",  "⚠ CWC物質",
         RGBColor(0x30, 0x0a, 0x0a), C_RED_LT),
        ("  ポリアクリル酸（Lubrizol）","US",   " 0.8%", "EAR99",  "—",
         RGBColor(0x08, 0x14, 0x2c), C_WHITE70),
        ("  ベンゾトリアゾール（上海）", "CN",   " 0.1%", "EAR99",  "中国産",
         RGBColor(0x20, 0x08, 0x08), C_RED_LT),
        ("  超純水 UPW（自社製造）",    "JP",   "68.3%", "N/A",    "—",
         RGBColor(0x06, 0x12, 0x22), C_WHITE70),
    ]
    row_h = 0.38
    for i, (name, country, pct, eccn, note, bg, text_c) in enumerate(bom_rows):
        y = 2.0 + i * row_h
        add_rect(slide, col1_x, y, 4.6, row_h - 0.03, bg)
        row_text = f"{name:<28} {country:<5} {pct:<7} {eccn:<9} {note}"
        add_textbox(slide, col1_x + 0.05, y + 0.05, 4.5, row_h - 0.08,
                    row_text, font_size=8.5, color=text_c)

    # EAR warning
    add_rect(slide, col1_x, 3.95, 4.6, 0.55, RGBColor(0x3b, 0x1a, 0x02))
    add_multiline(slide, col1_x + 0.1, 4.0, 4.4, 0.5, [
        ("🚩 EAR de minimis 超過",                         9, True,  C_WARN,    PP_ALIGN.LEFT),
        ("US原産品合計 25.8% > 25%閾値 → 中国向け再輸出に EAR ライセンス審査必要",
                                                           8.5, False, C_WHITE70, PP_ALIGN.LEFT),
    ])

    # ── Right: AI Result ──
    col2_x = 5.2

    # AI result box
    add_rect(slide, col2_x, 1.0, 4.5, 2.55, RGBColor(0x05, 0x10, 0x28))
    add_textbox(slide, col2_x + 0.1, 1.05, 4.3, 0.25,
                "🤖  AI HSコード判定結果（期待値）", font_size=9.5, bold=True, color=C_S2)

    hs_ranks = [
        ("1", "3824.99", "Chemical preparations NEC",    "1C350", C_S2_LIGHT, C_RED_LT),
        ("2", "2811.22", "Silicon dioxide（シリカ）",      "EAR99", C_WHITE70,  C_GRAY),
        ("3", "3824.40", "Prepared additives（調整剤）",   "EAR99", C_WHITE70,  C_GRAY),
        ("4", "3824.50", "Non-refractory mortars",        "EAR99", C_WHITE70,  C_GRAY),
        ("5", "2847.00", "Hydrogen peroxide（参考）",      "1C350", C_WHITE70,  C_RED_LT),
    ]
    for i, (rank, code, desc, eccn, tc, ec) in enumerate(hs_ranks):
        y = 1.35 + i * 0.43
        bg = RGBColor(0x0e, 0x24, 0x4a) if i == 0 else RGBColor(0x08, 0x16, 0x32)
        add_rect(slide, col2_x + 0.1, y, 4.2, 0.38, bg)
        add_multiline(slide, col2_x + 0.2, y + 0.04, 4.0, 0.32, [
            (f"#{rank}  {code}   {desc:<30}  [{eccn}]",
             9 if i == 0 else 8.5,
             i == 0, tc, PP_ALIGN.LEFT),
        ])

    # Apply note
    add_rect(slide, col2_x, 3.6, 4.5, 0.4, RGBColor(0x08, 0x22, 0x4e))
    add_textbox(slide, col2_x + 0.1, 3.66, 4.3, 0.3,
                "✅  「反映」ボタン → HS 3824.99 + ECCN 1C350 を品目フォームへ一括書込み",
                font_size=9, color=C_S2_LIGHT)

    # BOM CSV note
    add_rect(slide, col2_x, 4.1, 4.5, 0.5, RGBColor(0x06, 0x14, 0x2c))
    add_multiline(slide, col2_x + 0.1, 4.15, 4.3, 0.4, [
        ("📎  BOM CSVファイル（別添）",                          9,   True,  C_S2,     PP_ALIGN.LEFT),
        ("data/demo/cmp_slurry_bom.csv  —  7原料 全ECCN・HS記載", 8.5, False, C_WHITE70, PP_ALIGN.LEFT),
    ])

    # DAP alert
    add_rect(slide, col1_x, 4.75, 9.3, 0.42, RGBColor(0x0a, 0x18, 0x10))
    add_textbox(slide, col1_x + 0.1, 4.8, 9.0, 0.34,
                "👔 DAP: 「BOM-001 (コロイダルシリカ) が US 原産 25% です。EAR de minimis チェックを実施してください。」",
                font_size=9, color=C_WARN)

    # Footer
    add_rect(slide, 0, 5.5, 10, 0.12, C_S2)
    add_textbox(slide, 0.3, 5.5, 9.4, 0.12,
                "AI_TradeManagement  —  Demo Scenario 02  |  http://localhost:8002",
                font_size=7, color=C_WHITE)


def build_s3_slide(prs: Presentation):
    """シナリオ3: 半導体メモリ キャッチオール審査"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x02, 0x0e, 0x06))

    # Header
    add_rect(slide, 0, 0, 10, 0.85, RGBColor(0x06, 0x18, 0x0c))
    add_rect(slide, 0, 0.85, 10, 0.03, C_S3)

    add_multiline(slide, 0.3, 0.08, 7.5, 0.7, [
        ("SCENARIO  03", 9, True, C_S3, PP_ALIGN.LEFT),
        ("輸出審査 — 半導体メモリ キャッチオール確認 & 正規審査", 18, True, C_S3_LIGHT, PP_ALIGN.LEFT),
    ])
    add_textbox(slide, 7.5, 0.2, 2.3, 0.45,
                "ai_validation\n:8001",
                font_size=9, bold=True, color=C_S3, align=PP_ALIGN.RIGHT)

    # ── Left: Product + Catchall ──
    col1_x = 0.3

    # Product info
    add_rect(slide, col1_x, 1.0, 4.5, 0.85, RGBColor(0x06, 0x1a, 0x0c))
    add_textbox(slide, col1_x + 0.1, 1.05, 4.3, 0.22,
                "▌ 品目 & 取引情報", font_size=9, bold=True, color=C_S3)
    add_multiline(slide, col1_x + 0.15, 1.3, 4.2, 0.5, [
        ("品目: LPDDR5 16Gb DRAM チップ（7nm世代）  /  HS 854232",  9.5, False, C_WHITE70,  PP_ALIGN.LEFT),
        ("ECCN: 4A994 / 3A001  ｜  取引先: 深圳科技集団有限公司",    9.5, False, C_WHITE70,  PP_ALIGN.LEFT),
        ("仕向地: CN（新疆向け通信インフラプロジェクト）",              9.5, True,  C_RED_LT,  PP_ALIGN.LEFT),
    ])

    # Catchall 6steps
    add_rect(slide, col1_x, 1.95, 4.5, 0.28, RGBColor(0x0a, 0x28, 0x14))
    add_textbox(slide, col1_x + 0.1, 1.98, 4.3, 0.22,
                "⚖️  キャッチオール 6ステップ判定", font_size=9.5, bold=True, color=C_S3)

    steps = [
        ("1", "エンブレムト国チェック（E:1）",           "PASS",              RGBColor(0x06, 0x1a, 0x0e), C_S3_LIGHT, False),
        ("2", "ホワイト国チェック（A:1〜A:6）",         "FAIL: CN 非該当",    RGBColor(0x28, 0x08, 0x08), C_RED_LT,   True),
        ("3", "EAR Country Chart（AT/CC列）照合",       "TRIGGERED ⬅",       RGBColor(0x30, 0x0a, 0x0a), C_RED_LT,   True),
        ("4", "Red Flag 7項目（軍用転用・新疆）",        "FLAG ×3 ⚠",         RGBColor(0x2e, 0x18, 0x04), C_WARN,     True),
        ("5", "総合スコアリング",                        "Score: 78/100",     RGBColor(0x20, 0x0c, 0x0c), C_RED_LT,   True),
        ("6", "最終判定",                               "REQUIRES_PERMIT 🔴", RGBColor(0x38, 0x06, 0x06), C_RED_LT,   True),
    ]
    for i, (num, label, result, bg, tc, highlight) in enumerate(steps):
        y = 2.28 + i * 0.375
        add_rect(slide, col1_x, y, 4.5, 0.35, bg)
        add_multiline(slide, col1_x + 0.1, y + 0.05, 4.2, 0.27, [
            (f"Step {num}  {label:<40}  {result}", 8.5, highlight, tc, PP_ALIGN.LEFT),
        ])

    # ── Right: Pipeline + Result ──
    col2_x = 5.1

    # Pipeline
    add_rect(slide, col2_x, 1.0, 4.6, 0.28, RGBColor(0x06, 0x1a, 0x0c))
    add_textbox(slide, col2_x + 0.1, 1.03, 4.4, 0.22,
                "▌ AI 審査 Pipeline", font_size=9, bold=True, color=C_S3)

    pipe_steps = [
        ("usage\nexpand",     False),
        ("matrix_match\nLayer A",  True),
        ("dual_list\ncheck",       True),
        ("patent\nretrieve",       False),
        ("catchall\nassess",       True),
        ("finalize\nreport",       True),
    ]
    pipe_w = 0.72
    for i, (label, active) in enumerate(pipe_steps):
        x = col2_x + 0.05 + i * (pipe_w + 0.05)
        bg = RGBColor(0x0e, 0x30, 0x18) if active else RGBColor(0x06, 0x18, 0x0c)
        tc = C_S3_LIGHT if active else C_GRAY
        add_rect(slide, x, 1.3, pipe_w, 0.65, bg)
        add_textbox(slide, x + 0.03, 1.35, pipe_w - 0.06, 0.55,
                    label, font_size=8, color=tc, align=PP_ALIGN.CENTER, wrap=True)

    # Final result box
    add_rect(slide, col2_x, 2.05, 4.6, 1.65, RGBColor(0x2a, 0x06, 0x06))
    add_textbox(slide, col2_x + 0.2, 2.1, 4.2, 0.38,
                "🔴  REQUIRES_PERMIT",
                font_size=18, bold=True, color=C_RED_LT, align=PP_ALIGN.CENTER)
    add_multiline(slide, col2_x + 0.2, 2.52, 4.2, 1.1, [
        ("• 外為法 EL-3 該当 → 経済産業大臣許可申請必要",   9.5, False, C_WHITE70, PP_ALIGN.LEFT),
        ("• EAR AT列・CC列 トリガー → BIS ライセンス要",    9.5, False, C_WHITE70, PP_ALIGN.LEFT),
        ("• Red Flag: 軍用転用・新疆設置・MCF顧客 ×3件",   9.5, True,  C_WARN,    PP_ALIGN.LEFT),
        ("• 制裁リストスクリーニング: 実施必須（未完了警告）", 9.5, False, C_RED_LT,  PP_ALIGN.LEFT),
    ])

    # PDF report note
    add_rect(slide, col2_x, 3.8, 4.6, 0.65, RGBColor(0x06, 0x1e, 0x10))
    add_multiline(slide, col2_x + 0.15, 3.85, 4.3, 0.55, [
        ("📄  PDF レポート出力（4章構成）",                      9.5, True,  C_S3_LIGHT, PP_ALIGN.LEFT),
        ("Section 1: 品目概要  |  2: 外為法判定  |  3: デュアルリスト  |  4: キャッチオール詳細",
                                                               8.5, False, C_WHITE70,  PP_ALIGN.LEFT),
    ])

    # DAP alert
    add_rect(slide, col1_x, 4.75, 9.3, 0.42, RGBColor(0x18, 0x14, 0x04))
    add_textbox(slide, col1_x + 0.1, 4.8, 9.0, 0.34,
                "👔 DAP: 「本案件は仕向地 CN / 新疆関連プロジェクトで高リスクです。取引先スクリーニング (http://localhost:8005) を先に実施してください。」",
                font_size=9, color=C_WARN)

    # Footer
    add_rect(slide, 0, 5.5, 10, 0.12, C_S3)
    add_textbox(slide, 0.3, 5.5, 9.4, 0.12,
                "AI_TradeManagement  —  Demo Scenario 03  |  http://localhost:8001",
                font_size=7, color=C_BG)


# ─── PPTX Builder ─────────────────────────────────────────────────────────────

def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    build_title_slide(prs)
    build_s1_slide(prs)
    build_s2_slide(prs)
    build_s3_slide(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"✅  PPTX saved → {OUT_PATH}")
    return OUT_PATH


# ─── Google Drive Upload ───────────────────────────────────────────────────────

SCOPES = [
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/presentations",
]


def get_credentials() -> Credentials:
    creds = None
    if TOKEN_PATH.exists():
        creds = Credentials.from_authorized_user_file(str(TOKEN_PATH), SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
            TOKEN_PATH.write_text(creds.to_json())
        else:
            from google_auth_oauthlib.flow import InstalledAppFlow
            flow = InstalledAppFlow.from_client_secrets_file(str(CREDS_PATH), SCOPES)
            creds = flow.run_local_server(port=0)
            TOKEN_PATH.write_text(creds.to_json())
    return creds


def upload_to_google_slides(pptx_path: Path) -> str:
    """PPTX を Google Drive にアップロードし Google Slides に変換。URL を返す。"""
    creds = get_credentials()
    drive = build("drive", "v3", credentials=creds)

    file_metadata = {
        "name":     "AI_TradeManagement デモシナリオスライド",
        "mimeType": "application/vnd.google-apps.presentation",
    }
    media = MediaFileUpload(
        str(pptx_path),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=True,
    )

    print("⬆️   Google Drive にアップロード中...")
    file = drive.files().create(
        body=file_metadata,
        media_body=media,
        fields="id,webViewLink",
    ).execute()

    file_id   = file["id"]
    view_link = file.get("webViewLink", f"https://docs.google.com/presentation/d/{file_id}/edit")
    print(f"✅  Google Slides 作成完了！")
    print(f"    ファイル ID : {file_id}")
    print(f"    URL         : {view_link}")
    return view_link


# ─── Main ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("=" * 60)
    print("  AI_TradeManagement Demo Slides Builder")
    print("=" * 60)

    pptx_path = build_pptx()

    print()
    print("─" * 60)
    print("  Google Slides にアップロードします...")
    print("─" * 60)

    try:
        url = upload_to_google_slides(pptx_path)
        print()
        print("🎉  完了！ブラウザで開くには:")
        print(f"    open '{url}'")
    except Exception as e:
        print(f"⚠️   Google Slides アップロードに失敗しました: {e}")
        print(f"    ローカルの PPTX ファイルを直接開いてください: {pptx_path}")
        print(f"    open '{pptx_path}'")
